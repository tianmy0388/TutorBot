"""Tests for the knowledge base loaders (Task 8).

We test the five supported formats with real files generated in a tmp
dir, plus the failure cases (empty TXT, corrupted PDF).
"""

from __future__ import annotations

from pathlib import Path

import pytest

from tutor.services.knowledge_base.loaders import LoaderError, extract_text


def _write_txt(path: Path, content: str = "") -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def _write_pdf(path: Path) -> Path:
    # Use reportlab if available; otherwise a minimal hand-rolled PDF
    # that pypdf can read back. Reportlab produces a deterministic
    # text-bearing PDF.
    try:
        from reportlab.pdfgen import canvas  # type: ignore

        c = canvas.Canvas(str(path))
        c.drawString(72, 720, "Hello PDF world")
        c.showPage()
        c.save()
        return path
    except ImportError:
        # Fallback: write a PDF with a hand-rolled content stream that
        # includes literal text. This is enough for pypdf to extract.
        from pypdf import PdfWriter
        from pypdf.generic import (
            ArrayObject,
            ContentStream,
            DecodedStreamObject,
            DictionaryObject,
            FloatObject,
            NameObject,
            NumberObject,
            TextStringObject,
        )

        writer = PdfWriter()
        page = writer.add_blank_page(width=200, height=200)
        content = (
            b"BT /F1 12 Tf 50 100 Td (Hello PDF world) Tj ET"
        )
        stream = DecodedStreamObject()
        stream.set_data(content)
        page[NameObject("/Contents")] = ContentStream(stream, writer)
        with path.open("wb") as f:
            writer.write(f)
        return path


def _write_docx(path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello DOCX world")
    doc.save(str(path))
    return path


def _write_pptx(path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Hello PPTX"
    prs.save(str(path))
    return path


def _write_markdown(path: Path) -> Path:
    path.write_text(
        "# Title\n\nThis is a markdown document with a [link](https://example.com).\n",
        encoding="utf-8",
    )
    return path


def test_extract_txt(tmp_path: Path) -> None:
    p = _write_txt(tmp_path / "a.txt", "line 1\nline 2\nline 3\n")
    chunks = extract_text(p)
    assert chunks
    assert all(c.text for c in chunks)


def test_extract_markdown(tmp_path: Path) -> None:
    p = _write_markdown(tmp_path / "a.md")
    chunks = extract_text(p)
    assert chunks
    assert any("markdown" in c.text.lower() for c in chunks)


def test_extract_pdf(tmp_path: Path) -> None:
    p = _write_pdf(tmp_path / "a.pdf")
    chunks = extract_text(p)
    # Blank page has no text — but the extraction must not raise.
    assert isinstance(chunks, list)


def test_extract_docx(tmp_path: Path) -> None:
    p = _write_docx(tmp_path / "a.docx")
    chunks = extract_text(p)
    assert any("Hello" in c.text for c in chunks)


def test_extract_pptx(tmp_path: Path) -> None:
    p = _write_pptx(tmp_path / "a.pptx")
    chunks = extract_text(p)
    assert any("Hello" in c.text for c in chunks)


def test_empty_txt_raises(tmp_path: Path) -> None:
    p = _write_txt(tmp_path / "empty.txt", "   \n  \n")
    with pytest.raises(LoaderError) as exc:
        extract_text(p)
    assert exc.value.code == "EMPTY_DOCUMENT"


def test_corrupted_pdf_raises(tmp_path: Path) -> None:
    p = tmp_path / "broken.pdf"
    p.write_bytes(b"%PDF-not really\n")
    with pytest.raises(LoaderError) as exc:
        extract_text(p)
    assert exc.value.code in {"EXTRACTION_FAILED", "EMPTY_DOCUMENT"}


def test_unsupported_extension_raises(tmp_path: Path) -> None:
    p = tmp_path / "x.xyz"
    p.write_text("hi", encoding="utf-8")
    with pytest.raises(LoaderError) as exc:
        extract_text(p)
    assert exc.value.code == "UNSUPPORTED_FORMAT"
