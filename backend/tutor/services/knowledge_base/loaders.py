"""File loaders for knowledge base ingestion (Task 8).

Extracts plain text (with simple source anchors) from PDF, DOCX, PPTX,
Markdown and TXT. All loaders return a list of
:class:`ExtractedChunk` (text + anchor) so the downstream chunker can
split further if needed.

The loaders are deliberately best-effort: a corrupted PDF is allowed
to raise :class:`LoaderError` so the service can mark the document
``failed`` with a stable error code.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


class LoaderError(Exception):
    """A loader-specific failure (e.g. empty document, corrupt PDF)."""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


@dataclass
class ExtractedChunk:
    text: str
    anchor: str  # page / paragraph / slide


def extract_text(path: Path) -> list[ExtractedChunk]:
    """Dispatch to the right loader based on the file extension."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".md":
        return _extract_markdown(path)
    if ext == ".txt":
        return _extract_txt(path)
    raise LoaderError("UNSUPPORTED_FORMAT", f"unsupported: {ext}")


# ---------------------------------------------------------------------------
# Loaders
# ---------------------------------------------------------------------------

def _extract_txt(path: Path) -> list[ExtractedChunk]:
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = path.read_text(encoding="gbk", errors="replace")
    if not content.strip():
        raise LoaderError("EMPTY_DOCUMENT", "文件为空")
    # One chunk per ~30 lines.
    chunks: list[ExtractedChunk] = []
    lines = content.splitlines()
    for i in range(0, len(lines), 30):
        block = "\n".join(lines[i : i + 30]).strip()
        if block:
            chunks.append(ExtractedChunk(text=block, anchor=f"lines {i+1}-{min(i+30, len(lines))}"))
    if not chunks:
        raise LoaderError("EMPTY_DOCUMENT", "未提取到文本内容")
    return chunks


def _extract_markdown(path: Path) -> list[ExtractedChunk]:
    return _extract_txt(path)  # same shape


def _extract_pdf(path: Path) -> list[ExtractedChunk]:
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover
        raise LoaderError("DEPENDENCY_MISSING", f"pypdf missing: {e}") from e
    try:
        reader = PdfReader(str(path))
    except Exception as e:  # noqa: BLE001
        raise LoaderError("EXTRACTION_FAILED", f"PDF 读取失败: {e}") from e
    chunks: list[ExtractedChunk] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            raise LoaderError("EXTRACTION_FAILED", f"PDF 第 {i} 页失败: {e}") from e
        text = text.strip()
        if text:
            # Replace lone surrogates so the chunks json write
            # doesn't trip the strict utf-8 encoder later.
            try:
                text.encode("utf-8").decode("utf-8")
            except UnicodeEncodeError:
                text = text.encode("utf-8", errors="replace").decode(
                    "utf-8", errors="replace"
                )
            chunks.append(ExtractedChunk(text=text, anchor=f"page {i}"))
    # A PDF with no extractable text (image-only, blank, scanned with
    # no OCR) is a valid input — return an empty chunk list and let
    # the caller decide whether that's a failure (the service
    # currently treats it as EMPTY_DOCUMENT, but downstream code can
    # also choose to skip such files). Previously the loader raised
    # unconditionally, which made it impossible to distinguish a
    # "nothing to extract" page from a corrupt one in unit tests.
    return chunks


def _extract_docx(path: Path) -> list[ExtractedChunk]:
    try:
        from docx import Document
    except ImportError as e:  # pragma: no cover
        raise LoaderError("DEPENDENCY_MISSING", f"python-docx missing: {e}") from e
    try:
        doc = Document(str(path))
    except Exception as e:  # noqa: BLE001
        raise LoaderError("EXTRACTION_FAILED", f"DOCX 读取失败: {e}") from e
    chunks: list[ExtractedChunk] = []
    para_index = 0
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            para_index += 1
            chunks.append(ExtractedChunk(text=text, anchor=f"paragraph {para_index}"))
    if not chunks:
        raise LoaderError("EMPTY_DOCUMENT", "DOCX 无可提取文本")
    return chunks


def _extract_pptx(path: Path) -> list[ExtractedChunk]:
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise LoaderError("DEPENDENCY_MISSING", f"python-pptx missing: {e}") from e
    try:
        pres = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        raise LoaderError("EXTRACTION_FAILED", f"PPTX 读取失败: {e}") from e
    chunks: list[ExtractedChunk] = []
    for i, slide in enumerate(pres.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in p.runs).strip()
                    if text:
                        parts.append(text)
        text = "\n".join(parts).strip()
        if text:
            chunks.append(ExtractedChunk(text=text, anchor=f"slide {i}"))
    if not chunks:
        raise LoaderError("EMPTY_DOCUMENT", "PPTX 无可提取文本")
    return chunks


__all__ = [
    "ExtractedChunk",
    "LoaderError",
    "extract_text",
]
