"""PPT generation service — wraps ``python-pptx`` to turn Markdown
into a slide deck.

Design notes
------------

We don't try to be fancy. The slide structure is:

    Slide 1: cover — topic, subtitle, optional audience tag
    Slide 2: outline — bullets of every H2 in the source
    Slide 3..N: body — one slide per H2 with bullets from H3 + paragraph
    Last slide: summary — key takeaways (first 3-5 sentences of source)

Bullets are extracted with a tiny Markdown parser (line-based, supports
H2/H3, unordered lists, fenced code blocks as monospace, paragraphs as
notes). If a section is too long we cap at ~6 bullets per slide and
insert a "continued" slide if needed.

Output file is written to ``<data_dir>/ppt/<package_id>/<resource_id>.pptx``
and the path is stored in ``Resource.format_specific["pptx_path"]``.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from loguru import logger
from pptx import Presentation
from pptx.util import Inches, Pt

from tutor.services.config.settings import get_settings


# ---------------------------------------------------------------------------
# Markdown slice
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    """One slide worth of content."""

    title: str
    bullets: list[str] = None  # type: ignore[assignment]
    notes: str = ""

    def __post_init__(self) -> None:
        if self.bullets is None:
            self.bullets = []


_H2_RE = re.compile(r"^##\s+(.+?)\s*$")
_H3_RE = re.compile(r"^###\s+(.+?)\s*$")
_UL_RE = re.compile(r"^[-*]\s+(.+)$")
_CODE_FENCE_RE = re.compile(r"^```")
_TITLE_RE = re.compile(r"^#\s+(.+?)\s*$")


def _strip_fences(text: str) -> str:
    """Remove ``` fences but keep the code content as preformatted."""
    out: list[str] = []
    in_fence = False
    for line in text.splitlines():
        if _CODE_FENCE_RE.match(line):
            in_fence = not in_fence
            out.append("")  # gap
            continue
        out.append(line)
    return "\n".join(out)


def slice_markdown_to_slides(
    topic: str, markdown: str, *, max_bullets_per_slide: int = 6
) -> list[Slide]:
    """Convert Markdown into a structured list of :class:`Slide`.

    Always returns at least one slide (a cover).
    """
    text = _strip_fences(markdown or "")

    title = topic
    sections: list[tuple[str, list[str]]] = []  # (h2 title, [bullets])
    current_h2: str | None = None
    current_bullets: list[str] = []

    for line in text.splitlines():
        line = line.rstrip()
        if not line:
            continue
        m = _TITLE_RE.match(line)
        if m:
            title = m.group(1).strip()
            continue
        m = _H2_RE.match(line)
        if m:
            if current_h2 is not None:
                sections.append((current_h2, current_bullets))
            current_h2 = m.group(1).strip()
            current_bullets = []
            continue
        m = _H3_RE.match(line)
        if m:
            current_bullets.append(f"— {m.group(1).strip()}")
            continue
        m = _UL_RE.match(line)
        if m:
            current_bullets.append(m.group(1).strip())
            continue
        # Plain paragraph — only add as bullet if we have a current section
        # AND the section is empty (so a paragraph doesn't drown the slide)
        if current_h2 is not None and not current_bullets:
            current_bullets.append(line.strip())

    if current_h2 is not None:
        sections.append((current_h2, current_bullets))

    slides: list[Slide] = [
        Slide(title=title, bullets=[], notes="课程概览"),
        Slide(title="本节目录", bullets=[s[0] for s in sections[:8]] or [topic]),
    ]

    for h2, bullets in sections:
        if not bullets:
            slides.append(Slide(title=h2, bullets=[], notes=""))
            continue
        # Chunk bullets into slides of `max_bullets_per_slide`
        for i in range(0, len(bullets), max_bullets_per_slide):
            chunk = bullets[i : i + max_bullets_per_slide]
            label = h2
            if i > 0:
                label = f"{h2} (续)"
            slides.append(Slide(title=label, bullets=chunk))

    # Summary slide from the first ~3 sentences of the source if available
    summary = _first_n_sentences(text, n=5)
    if summary:
        slides.append(Slide(title="关键要点", bullets=summary))

    return slides


_SENTENCE_RE = re.compile(r"(?<=[。！？!?\.])\s")


def _first_n_sentences(text: str, n: int = 5) -> list[str]:
    """Pick the first n sentences from prose (Chinese + English punctuation)."""
    # Strip headings and code-fences-ish whitespace
    plain_lines = [
        ln.strip() for ln in text.splitlines()
        if ln.strip() and not ln.startswith("#") and not ln.startswith(("-", "*", "```"))
    ]
    plain = " ".join(plain_lines)
    if not plain:
        return []
    # Split on sentence terminators (mixed CJK + ASCII)
    parts = re.split(r"(?<=[。！？!?\.])", plain)
    parts = [p.strip() for p in parts if p.strip()]
    return parts[:n]


# ---------------------------------------------------------------------------
# python-pptx rendering
# ---------------------------------------------------------------------------


def _set_bullets(slide, bullets: Iterable[str]) -> None:
    body = slide.shapes.placeholders[1].text_frame if len(slide.shapes.placeholders) > 1 else None
    if body is None:
        return
    body.clear()
    first = True
    for b in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)


def render_slides(slides: list[Slide], output_path: Path, *, title: str = "") -> Path:
    """Write a :class:`Presentation` to ``output_path``.

    Uses the default python-pptx title-and-content layout for body
    slides. Falls back gracefully when text is empty.
    """
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]   # Title Slide
    bullet_layout = prs.slide_layouts[1]  # Title and Content

    for i, s in enumerate(slides):
        layout = bullet_layout if (i > 0 or s.bullets) else title_layout
        slide = prs.slides.add_slide(layout)
        # Title placeholder
        try:
            slide.shapes.title.text = s.title or title or "未命名"
        except Exception:
            pass
        # Body bullets (skip for cover slide)
        if s.bullets and i > 0:
            _set_bullets(slide, s.bullets)
        elif not s.bullets and i == 0:
            # Subtitle hint on the cover
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = title or ""

    prs.save(str(output_path))
    logger.info(f"PPT written to {output_path} ({len(slides)} slides)")
    return output_path


# ---------------------------------------------------------------------------
# High-level facade
# ---------------------------------------------------------------------------


class PPTGenerationService:
    """Stateless facade for ``python-pptx`` rendering.

    Use :meth:`build` to slice a Markdown source into slides and write
    the resulting deck to disk. Returns the path.
    """

    def __init__(self, output_dir: Path | None = None) -> None:
        self.output_dir = output_dir or (get_settings().data_dir / "ppt")

    def build(
        self,
        *,
        topic: str,
        markdown: str,
        package_id: str,
        resource_id: str,
        title: str | None = None,
    ) -> Path:
        slides = slice_markdown_to_slides(topic, markdown)
        out_dir = self.output_dir / package_id
        out_path = out_dir / f"{resource_id}.pptx"
        render_slides(slides, out_path, title=title or topic)
        return out_path


# Singleton accessor
_service: PPTGenerationService | None = None


def get_ppt_service() -> PPTGenerationService:
    global _service
    if _service is None:
        _service = PPTGenerationService()
    return _service


__all__ = [
    "PPTGenerationService",
    "Slide",
    "get_ppt_service",
    "render_slides",
    "slice_markdown_to_slides",
]