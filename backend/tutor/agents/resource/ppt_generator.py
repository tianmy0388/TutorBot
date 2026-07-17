"""PPTGeneratorAgent — turns a document-style Resource into a .pptx deck.

Pipeline contract:

    topic: str
    source_content: str  (markdown; typically the ContentExpert output)
    profile: dict (optional)

    → Resource(type=PPT, format_specific={slide_count, slide_titles},
      confidence=~0.8)

The agent is intentionally **deterministic** when no LLM is available:
it slices the Markdown into slides via
:func:`tutor.services.ppt.slice_markdown_to_slides` and renders
straight away. With an LLM it could rewrite each section into tighter
slide bullets, but the deterministic path is enough for the MVP — the
output is already a usable teaching deck.
"""

from __future__ import annotations

import asyncio
import os
import threading
from contextlib import suppress
from typing import Any

from loguru import logger

from tutor.agents.base_agent import BaseAgent
from tutor.core.redaction import public_failure
from tutor.core.stream_bus import StreamBus
from tutor.services.ppt import get_ppt_service
from tutor.services.resource_package.schema import (
    PPTResource,
    Resource,
    ResourceType,
    build_resource,
)


class PPTGeneratorAgent(BaseAgent):
    """Convert a Markdown source into a PPT deck Resource."""

    agent_name = "ppt_generator"

    def __init__(
        self,
        ppt_service: Any | None = None,
        *,
        estimated_minutes: int = 12,
    ) -> None:
        super().__init__()
        self.ppt_service = ppt_service or get_ppt_service()
        self.estimated_minutes = estimated_minutes

    async def process(  # type: ignore[override]
        self,
        *,
        topic: str,
        source_content: str,
        profile: dict[str, Any] | None = None,
        package_id: str | None = None,
        stream: StreamBus | None = None,
    ) -> Resource:
        # Build a stable resource_id up front so we can also use it as
        # the on-disk filename.
        resource = build_resource(
            type=ResourceType.PPT,
            title=f"{topic} — PPT 教案",
            content=source_content or topic,
            difficulty=2,
            estimated_minutes=self.estimated_minutes,
            topic=topic,
            tags=["ppt", "教案", "outline"],
            confidence_score=0.78,
            generated_by=[self.agent_name],
            metadata={
                "source_chars": len(source_content or ""),
                "agent": self.agent_name,
            },
        )
        # Render (off-thread to avoid blocking the event loop on slow disks)
        pkg_id = package_id or "ad_hoc"
        cancel_event = threading.Event()
        publish_lock = threading.Lock()
        worker = asyncio.create_task(
            asyncio.to_thread(
                self.ppt_service.build,
                topic=topic,
                markdown=source_content or "",
                package_id=pkg_id,
                resource_id=resource.resource_id,
                title=topic,
                cancel_event=cancel_event,
                publish_lock=publish_lock,
            )
        )
        try:
            try:
                pptx_path = await asyncio.shield(worker)
            except Exception:  # noqa: BLE001
                failure = public_failure(
                    "PPT_RENDER_FAILED", "PPT rendering failed", retryable=True
                )
                logger.error("PPT_RENDER_FAILED agent={}", self.agent_name)
                if stream is not None:
                    await stream.error(
                        "PPT rendering failed",
                        source=self.agent_name,
                        metadata=failure,
                    )
                # A failed renderer did not create a usable educational resource.
                # Keep the typed diagnostic for trace/retry UI, but never retain
                # the source document in this failed artifact.
                resource.content = "PPT rendering failed. Please retry."
                resource.format_specific = {
                    "slide_count": 0,
                    "pptx_path": None,
                    "slide_titles": [],
                    "failure": failure,
                }
                resource.confidence_score = 0.0
                return resource

            # Populate format_specific from the on-disk artifact.
            try:
                slide_titles, slide_count = _peek_pptx(pptx_path)
            except Exception:  # noqa: BLE001
                logger.warning("PPT_INSPECTION_FAILED file_usable=true")
                slide_titles, slide_count = [], 0

            payload = PPTResource(
                slide_count=slide_count,
                pptx_path=str(pptx_path),
                slide_titles=slide_titles,
            )
            resource.format_specific = payload.model_dump()
            resource.metadata["pptx_filename"] = os.path.basename(str(pptx_path))
            resource.metadata["file_size"] = pptx_path.stat().st_size

            if stream is not None:
                await stream.observation(
                    f"PPT 已生成 ({slide_count} 张): {pptx_path.name}",
                    source=self.agent_name,
                    metadata={
                        "slide_count": slide_count,
                        "package_id": pkg_id,
                    },
                )
            return resource
        except asyncio.CancelledError:
            await _set_cancelled(cancel_event, publish_lock)
            with suppress(Exception):
                await asyncio.shield(worker)
            cleanup = getattr(self.ppt_service, "cleanup_cancelled", None)
            if cleanup is not None:
                cleanup(
                    package_id=pkg_id,
                    resource_id=resource.resource_id,
                    publish_lock=publish_lock,
                )
            raise


async def _set_cancelled(
    cancel_event: threading.Event,
    publish_lock: threading.Lock,
) -> None:
    """Set cancellation in publish order without blocking the event loop."""

    while not publish_lock.acquire(blocking=False):
        await asyncio.sleep(0)
    try:
        cancel_event.set()
    finally:
        publish_lock.release()


def _peek_pptx(path) -> tuple[list[str], int]:
    """Open a rendered pptx and return (slide_titles, count). Cheap."""
    from pptx import Presentation  # local import keeps module import lean

    prs = Presentation(str(path))
    titles: list[str] = []
    for slide in prs.slides:
        try:
            t = slide.shapes.title.text if slide.shapes.title else ""
        except Exception:
            t = ""
        titles.append(t.strip() or "(untitled)")
    return titles, len(prs.slides)


__all__ = ["PPTGeneratorAgent"]
